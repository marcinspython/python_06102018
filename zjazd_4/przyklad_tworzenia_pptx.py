import csv
from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title

body_shape = shapes.placeholders[1]

title_shape.text = "Jakiś tekst"

tf = body_shape.text_frame
tf.text = "Zawartość tekst frame"

# with open("report.csv") as f:
#     data = csv.reader()
p = tf.add_paragraph()
p.text = "Kobiety"
p.level = 1

p = tf.add_paragraph()
p.text = "Przeżyło"
p.level = 2

prs.save('raport.pptx')

